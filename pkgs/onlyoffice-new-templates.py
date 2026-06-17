# Generates valid *blank* office documents used as OnlyOffice "new file"
# templates (newFileTemplate/<locale>/new.<ext>). A 0-byte placeholder is not a
# valid OOXML container, so OnlyOffice's x2t conversion fails with ExitCode 80
# ("an error has occurred while opening the file") when seeding a new document.
# Only the creatable types (app_registry allow_creation) need real content.
#
# Run as: python3 onlyoffice-new-templates.py <output-dir>
import sys

out = sys.argv[1]

from docx import Document
Document().save(out + "/new.docx")

import openpyxl
openpyxl.Workbook().save(out + "/new.xlsx")

from pptx import Presentation
pres = Presentation()
pres.slides.add_slide(pres.slide_layouts[6])  # index 6 = "Blank" layout
pres.save(out + "/new.pptx")

from reportlab.pdfgen import canvas
pdf = canvas.Canvas(out + "/new.pdf")
pdf.showPage()
pdf.save()
